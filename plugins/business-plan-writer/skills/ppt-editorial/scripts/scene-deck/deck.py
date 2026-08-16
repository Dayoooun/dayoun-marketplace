# -*- coding: utf-8 -*-
"""씬 덱 단일 진입점 (2026-08-01)

## 왜 필요한가
프리셋·레이아웃·폰트·사진·수정 모듈이 다 갖춰졌는데, 정작 덱을 만들 때마다
`build.py` 보일러플레이트(스타일 조립 → jobs.json → 생성 → 레이아웃 세팅 → PDF)를
새로 써야 했다. 실측으로 6개 프로젝트가 같은 코드를 반복했고, 그때마다 편차가 생겼다.

## 사용법
```python
from deck import Deck

d = Deck(domain="식음료", foot="OO식품", title="사업 소개")
d.slide("L", "THE PROBLEM", ["산지는 좋은데", "팔 길이 없습니다"],
        ["좋은 원물을 확보해도 판로가 없으면 재고가 된다"],
        scene="a lone food stall on an isolated platform")
d.slide("W", "OUR PROCESS", ["산지에서 식탁까지"], ["4단계 일괄 관리"],
        scene="four pedestals in a row: crate, machine, package, delivery box",
        labels=["수급", "가공", "포장", "유통"])
d.slide("L", "TRACK RECORD", ["숫자로 증명합니다"], ["공공기관 검증"],
        scene="a tall stack of layered plates with check badges",
        num=["2,800", "회", "연간 커밋"], chips=["HACCP", "직거래"])
d.photos(["a.jpg", "b.jpg", "c.jpg"], "ON SITE", ["현장에서 함께합니다"], ["3회차 밀착"])

d.generate(APPROVAL_DIGEST, APPROVAL_STORE)  # 승인된 씬만 생성
d.build(APPROVAL_DIGEST, APPROVAL_STORE)     # 조립 + PDF
```

한 번 쓴 덱은 `spec.json`으로 저장되어 `revise.py`로 수정할 수 있다.
"""
import os, sys, json, subprocess, glob
from pathlib import Path

HERE = os.path.dirname(os.path.abspath(__file__))
SCRIPTS = os.path.dirname(HERE)
if not os.path.exists(os.path.join(SCRIPTS, "codex_parallel_gen.py")):
    SCRIPTS = HERE  # git 추적용 평면 harness 사본
for path in (HERE, SCRIPTS):
    if path not in sys.path:
        sys.path.insert(0, path)

from presets import preset, style_block          # noqa: E402
import layout_engine as LE                        # noqa: E402
from codex_parallel_gen import scene_safe_zone_receipt  # noqa: E402
from approved_inputs import (  # noqa: E402
    ApprovalError,
    SCENE_RENDERER_VERSION,
    digest_value,
    resolve_approved_bundle,
    resolve_value,
    sha256_file,
)

GEN = os.path.join(SCRIPTS, "codex_parallel_gen.py")

SAFE = ("\n★★ SAFE ZONES — OVERRIDES EVERYTHING:\n"
        "1) Keep TOP, BOTTOM, LEFT, and RIGHT 18% completely empty.\n"
        "2) Every screen, card, button, connector, label, shadow, and module must fit inside the central 64%.\n"
        "3) Scale down and center the entire object group; no touching, cropping, or edge exits.\n")

TAIL = ("\nUse ONLY the image-generation capability - ABSOLUTELY NO Python/PIL/code drawing.\n"
        "Korean text inside the image must render perfectly — no '?', no broken glyphs.\n"
        "결과를 반드시 __out.png 로 저장하고 크기를 출력하라. PIL 드로잉 금지, 이미지 생성만.\n")

# 구도별 권장 씬 비율
RATIO = {"W": "16:9 wide", "F": "16:9 wide"}
DEFAULT_RATIO = "4:3 landscape"


def _safe_receipt_valid(path, margin=0.18):
    sidecar = Path(str(path) + ".safe.json")
    try:
        receipt = json.loads(sidecar.read_text(encoding="utf-8"))
        current = scene_safe_zone_receipt(path, margin)
    except (OSError, ValueError, TypeError, json.JSONDecodeError):
        return False
    return receipt == current


def _load_safe_receipt(path, margin=0.18):
    if not _safe_receipt_valid(path, margin):
        raise ApprovalError(f"scene safe-zone receipt is missing or stale: {path}")
    return json.loads(Path(str(path) + ".safe.json").read_text(encoding="utf-8"))


def _prepare_slide_outputs(directory, count):
    root = Path(directory)
    root.mkdir(parents=True, exist_ok=True)
    expected = {f"slide_{index:02d}.png" for index in range(1, count + 1)}
    for path in root.glob("slide_*.png"):
        generated_index = path.stem.removeprefix("slide_")
        if generated_index.isdigit() and path.name not in expected:
            path.unlink()


class Deck:
    def __init__(self, domain="it", foot="", title="", out_dir=None, brief=None):
        self.domain = domain
        self.preset = preset(domain)
        self.foot = foot
        self.title = title
        self.dir = os.path.abspath(out_dir or "deck_out")
        self.slides = []
        self._photo_refs = {}
        self.brief = None
        if brief is not None:
            from intake import require_confirmed
            self.brief = require_confirmed(brief)
        os.makedirs(self.dir, exist_ok=True)

    # ══════ 슬라이드 추가 ══════
    def slide(self, lay, eyebrow, head, sub, scene=None, labels=None,
              num=None, chips=None, items=None, ratio=None):
        """슬라이드 1장 추가.

        lay      L/S/W/C/A/F/T
        scene    씬 내용(영문 서술). None이면 텍스트만.
        labels   씬 안에 박을 한글 라벨
        num      ["2,800", "회", "연간 커밋"] 대형 수치
        chips    ["키워드1", "키워드2"] 칩
        items    [("제목","설명"), ...] T 구도 전용
        """
        sid = "s%02d" % (len(self.slides) + 1)
        s = {"_sid": len(self.slides) + 1, "lay": lay, "eyebrow": eyebrow,
             "scene": sid, "head": list(head), "sub": list(sub)}
        if scene:
            s["scene_body"] = scene
            s["scene_ratio"] = ratio or RATIO.get(lay, DEFAULT_RATIO)
        if labels:
            s["scene_labels"] = list(labels)
        for k, v in [("num", num), ("chips", chips), ("items", items)]:
            if v:
                s[k] = v
        self.slides.append(s)
        self.save()          # 매 추가마다 저장 — generate 후 load 가능하게(실측 버그)
        return self


    # ══════ 실무 필수 슬라이드 ══════
    def cover(self, eyebrow, head, sub=None, issuer=None, meta=None, scene=None):
        """표지. 크롬(쪽번호·푸터)이 붙지 않는다."""
        self.slide("COVER", eyebrow, head, sub or [], scene=scene,
                   ratio="4:3 landscape")
        self.slides[-1].update({"issuer": issuer or self.foot,
                                "meta": meta or [], "_nochrome": True})
        self.save()
        return self

    def agenda(self, items, head=None, eyebrow="AGENDA", scene=None):
        """목차. items=["배경","솔루션",...] 또는 [("01","배경"),...]"""
        norm = [(("%02d" % i), x) if isinstance(x, str) else tuple(x)
                for i, x in enumerate(items, 1)]
        self.slide("AGENDA", eyebrow, head or ["목차"], [], scene=scene)
        self.slides[-1]["items"] = norm
        self.save()
        return self

    def closing(self, head, sub=None, issuer=None, scene=None, eyebrow="THANK YOU"):
        """클로징."""
        self.slide("CLOSING", eyebrow, head, sub or [], scene=scene)
        self.slides[-1]["issuer"] = issuer or self.foot
        self.save()
        return self

    def photos(self, paths, eyebrow, head, sub, lay=None, labels=None, **kw):
        """실사진 슬라이드. 철칙 D — refs로 codex에 전달해 씬 안에서 함께 그린다."""
        from photos import plan, prep, prompt_block
        ready = prep(paths, os.path.join(self.dir, "photos_ready"))
        pl = plan(ready)
        lay = lay or {"hero": "L", "compare": "C", "sequence": "W", "grid": "F"}[pl["mode"]]
        self.slide(lay, eyebrow, head, sub,
                   scene="a %s composition integrating the provided photographs" % pl["mode"],
                   labels=labels, **kw)
        sid = self.slides[-1]["scene"]
        self._photo_refs[sid] = (ready, pl, labels)
        self.slides[-1]["_photo"] = [os.path.relpath(x, self.dir) for x in ready]
        self.save()
        return self

    # ══════ 정보 구도 — 씬 없이 텍스트만 렌더 ══════
    # 읽는 사람이 혼자 판단해야 하는 덱(교육·매뉴얼·심사)에서 쓴다.
    # 씬을 만들지 않으므로 생성 비용 0, 텍스트만 고치면 즉시 반영된다.
    # 정렬 규칙과 상세 인자는 info_layouts.py docstring 참조.

    def table(self, eyebrow, head, sub, columns, rows, **kw):
        """표. widths=[열 비율], align=[열별 정렬], source=출처."""
        import info_layouts as IL
        IL.register()
        return IL.table(self, eyebrow, head, sub, columns, rows, **kw)

    def example(self, eyebrow, head, sub, blocks, **kw):
        """대비. blocks=[(캡션, 본문, 해설), ...] — 첫 블록이 붉은 캡션."""
        import info_layouts as IL
        IL.register()
        return IL.example(self, eyebrow, head, sub, blocks, **kw)

    def matrix(self, eyebrow, head, sub, quadrants, x_label, y_label, **kw):
        """2×2. quadrants 는 좌상·우상·좌하·우하 순."""
        import info_layouts as IL
        IL.register()
        return IL.matrix(self, eyebrow, head, sub, quadrants, x_label, y_label, **kw)

    def bar(self, eyebrow, head, sub, bars, **kw):
        """가로 막대. bars=[(라벨, 값), ...]. 실측/예시를 source 에 밝힌다."""
        import info_layouts as IL
        IL.register()
        return IL.bar(self, eyebrow, head, sub, bars, **kw)

    def flow(self, eyebrow, head, sub, modules, **kw):
        """순서형 모듈 그래프. modules=[(step, title, detail), ...]."""
        import info_layouts as IL
        IL.register()
        return IL.flow(self, eyebrow, head, sub, modules, **kw)

    def genealogy(self, eyebrow, head, sub, modules, **kw):
        """개념 계보. modules=[(name, description, group), ...]."""
        import info_layouts as IL
        IL.register()
        return IL.genealogy(self, eyebrow, head, sub, modules, **kw)

    def prompt(self, eyebrow, head, sub, prompt_lines, **kw):
        """복사용 프롬프트 코드 블록."""
        import info_layouts as IL
        IL.register()
        return IL.prompt(self, eyebrow, head, sub, prompt_lines, **kw)

    # ══════ 씬 생성 ══════
    def jobs(self, only_missing=True):
        """씬 생성 잡 목록. 이미 있는 씬은 건너뛴다."""
        from photos import prompt_block
        scn = os.path.join(self.dir, "scenes")
        os.makedirs(scn, exist_ok=True)
        out = []
        for s in self.slides:
            if not s.get("scene_body"):
                continue
            sid = s["scene"]
            path = os.path.join(scn, "%s.png" % sid)
            if (
                only_missing
                and os.path.exists(path)
                and os.path.getsize(path) > 100_000
                and _safe_receipt_valid(path)
            ):
                continue
            p = style_block(self.domain) + "\nCOMPOSITION: " + s["scene_body"]
            if sid in self._photo_refs:
                ready, pl, labels = self._photo_refs[sid]
                p += prompt_block(pl, labels)
                refs = [os.path.relpath(x, self.dir) for x in ready]
            else:
                refs = []
                if s.get("scene_labels"):
                    p += ("\n\nTEXT LABELS rendered in the image "
                          "(bold Korean gothic, ink colour, small):\n  "
                          + "  ".join('"%s"' % x for x in s["scene_labels"]))
            p += SAFE + "\nOutput %s framing.\n" % s.get("scene_ratio", DEFAULT_RATIO) + TAIL
            out.append(
                {
                    "label": sid,
                    "refs": refs,
                    "out": os.path.join("scenes", "%s.png" % sid),
                    "prompt": p,
                    "safe_zone": 0.18,
                    "safe_frame": True,
                }
            )
        return out

    def approval_config(self):
        assets = {}
        for slide in self.slides:
            for relative in slide.get("_photo", []):
                path = (Path(self.dir) / relative).resolve()
                if not path.is_file():
                    raise ApprovalError(f"scene-deck referenced asset is missing: {relative}")
                assets[relative.replace("\\", "/")] = sha256_file(path)
        return {
            "domain": self.domain,
            "foot": self.foot,
            "title": self.title,
            "preset": self.preset,
            "assetDigests": assets,
        }

    def _scene_receipt_path(self):
        return os.path.join(self.dir, "scene-receipt.json")

    def _scene_job_record(self, job):
        output = Path(self.dir) / job["out"]
        return {
            "sceneId": job["label"],
            "promptDigest": digest_value(job["prompt"]),
            "referenceDigests": [
                sha256_file((Path(self.dir) / relative).resolve())
                for relative in job.get("refs", [])
            ],
            "outputDigest": sha256_file(output),
            "safeZone": _load_safe_receipt(output, float(job["safe_zone"])),
        }

    def _verify_scene_receipt(self, approval_digest, all_jobs):
        receipt_path = Path(self._scene_receipt_path())
        if not receipt_path.is_file():
            raise ApprovalError("scene receipt is missing")
        receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
        if (
            receipt.get("receiptVersion") != "dayoun-scene-render-receipt-v2"
            or receipt.get("approvalEnvelopeDigest") != approval_digest
            or receipt.get("rendererVersion") != SCENE_RENDERER_VERSION
        ):
            raise ApprovalError("scene receipt is stale")
        expected = [self._scene_job_record(job) for job in all_jobs]
        if receipt.get("scenes") != expected:
            raise ApprovalError("scene prompt, reference, or output changed after rendering")

    def _output_receipt_path(self):
        return os.path.join(self.dir, "deck-output-receipt.json")

    def _output_records(self):
        files = sorted(Path(self.dir, "out").glob("slide_*.png"))
        return [
            {"fileName": path.name, "sha256": sha256_file(path)}
            for path in files
        ]

    def _verify_output_receipt(self, approval_digest):
        path = Path(self._output_receipt_path())
        if not path.is_file():
            raise ApprovalError("deck output receipt is missing")
        receipt = json.loads(path.read_text(encoding="utf-8"))
        if (
            receipt.get("approvalEnvelopeDigest") != approval_digest
            or receipt.get("rendererVersion") != SCENE_RENDERER_VERSION
            or receipt.get("slides") != self._output_records()
        ):
            raise ApprovalError("assembled slide pixels changed after receipt")
    def _require_approved(self, approval_digest, approval_store):
        if not approval_digest or not approval_store:
            raise ApprovalError("scene-deck requires an approval digest and content-addressed store")
        envelope = resolve_value(approval_store, approval_digest)
        if envelope.get("rendererMode") != "scene-deck":
            raise ApprovalError("approval envelope does not authorize scene-deck")
        _payload, briefs, _manifest = resolve_approved_bundle(
            approval_store,
            envelope,
            expected_mode="scene-deck",
            expected_renderer_version=SCENE_RENDERER_VERSION,
        )
        approved_slides = [
            slide
            for brief in briefs
            for slide in brief.get("sceneDeckSlides", [])
        ]
        if not approved_slides:
            raise ApprovalError("approved scene-deck brief has no sceneDeckSlides")
        if digest_value(approved_slides) != digest_value(self.slides):
            raise ApprovalError("scene-deck render spec does not match the approved ordered briefs")
        approved_configs = [
            brief.get("sceneDeckConfig")
            for brief in briefs
            if brief.get("sceneDeckConfig") is not None
        ]
        if len(approved_configs) != 1:
            raise ApprovalError("approved scene-deck brief must contain one sceneDeckConfig")
        if digest_value(approved_configs[0]) != digest_value(self.approval_config()):
            raise ApprovalError("scene-deck renderer config or referenced asset changed")
        self._approval_digest = approval_digest
        return envelope

    def generate(self, approval_digest, approval_store, effort="high", timeout=800, loop=2):
        """Generate scenes only from a current digest-bound approval bundle."""
        self._require_approved(approval_digest, approval_store)
        all_jobs = self.jobs(only_missing=False)
        existing = [
            job
            for job in all_jobs
            if os.path.exists(os.path.join(self.dir, job["out"]))
        ]
        if existing:
            self._verify_scene_receipt(approval_digest, all_jobs)
        jobs = self.jobs()
        if not jobs:
            print("[생성] 승인 receipt와 모든 씬 digest 일치 — 생략")
            return self
        jobs_path = os.path.join(self.dir, "jobs.json")
        with open(jobs_path, "w", encoding="utf-8") as handle:
            json.dump(jobs, handle, ensure_ascii=False, indent=1)
        print("[생성] %d장 (effort=%s)" % (len(jobs), effort))
        completed = subprocess.run(
            [
                sys.executable,
                GEN,
                "jobs.json",
                "--retry",
                "1",
                "--loop",
                str(loop),
                "--timeout",
                str(timeout),
                "--effort",
                effort,
            ],
            cwd=self.dir,
            check=False,
        )
        failed = [
            job["label"]
            for job in all_jobs
            if not os.path.exists(os.path.join(self.dir, job["out"]))
            or os.path.getsize(os.path.join(self.dir, job["out"])) < 100_000
        ]
        if completed.returncode != 0 or failed:
            raise ApprovalError(
                "scene generation failed or produced undersized output: "
                + ",".join(failed or ["provider-exit"])
            )
        receipt = {
            "schemaVersion": "1.0.0",
            "receiptVersion": "dayoun-scene-render-receipt-v2",
            "approvalEnvelopeDigest": approval_digest,
            "rendererVersion": SCENE_RENDERER_VERSION,
            "scenes": [self._scene_job_record(job) for job in all_jobs],
        }
        Path(self._scene_receipt_path()).write_text(
            json.dumps(receipt, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
        self.save()
        return self

    # ══════ 조립 ══════
    def _apply_style(self):
        p = self.preset["palette"]
        LE.FAMILY = self.preset["fonts"]["head"]
        LE.SCN = os.path.join(self.dir, "scenes")
        LE.BLUE, LE.INK, LE.GREY, LE.LINE = p["hero"], p["ink"], p["grey"], p["line"]
        LE.FOOT = self.foot
        # 정보 구도가 스펙에 있으면 LAY 에 얹는다 — revise/재로드 경로에서도 안전하게
        if any(
            s["lay"] in ("TABLE", "EXAMPLE", "MATRIX", "BAR", "FLOW", "GENEALOGY", "PROMPT")
            for s in self.slides
        ):
            import info_layouts as IL
            IL.register()

    def build(self, approval_digest, approval_store, pdf=True, pptx=False):
        """Assemble outputs only from a current digest-bound approval bundle."""
        self._require_approved(approval_digest, approval_store)
        all_jobs = self.jobs(only_missing=False)
        missing = [
            job["label"]
            for job in all_jobs
            if not os.path.isfile(os.path.join(self.dir, job["out"]))
        ]
        if missing:
            raise ApprovalError("scene outputs are missing: " + ",".join(missing))
        if all_jobs:
            self._verify_scene_receipt(approval_digest, all_jobs)
        from PIL import Image, ImageDraw
        self._apply_style()
        out = os.path.join(self.dir, "out")
        os.makedirs(out, exist_ok=True)
        n = len(self.slides)
        _prepare_slide_outputs(out, n)
        for i, s in enumerate(self.slides, 1):
            im = Image.new("RGB", (LE.W, LE.H), LE.WHITE)
            d = ImageDraw.Draw(im)
            LE.LAY[s["lay"]](im, d, s)
            if not s.get("_nochrome"):
                LE.chrome(im, d, s["eyebrow"], "%02d" % i, n)
            im.save(os.path.join(out, "slide_%02d.png" % i))
        print("[조립] %d장 (도메인=%s 폰트=%s)" % (n, self.preset["key"], LE.FAMILY))
        output_receipt = {
            "schemaVersion": "1.0.0",
            "approvalEnvelopeDigest": approval_digest,
            "rendererVersion": SCENE_RENDERER_VERSION,
            "slides": self._output_records(),
        }
        Path(self._output_receipt_path()).write_text(
            json.dumps(output_receipt, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
        self.save()
        outs = []
        if pdf:
            outs.append(self.pdf(approval_digest, approval_store))
        if pptx:
            outs.append(self.pptx(approval_digest, approval_store))
        return outs[0] if len(outs) == 1 else (outs or None)

    def pptx(self, approval_digest, approval_store, name=None):
        """편집 가능한 PPTX 출력. 고객이 PPT 파일을 원할 때.
        슬라이드 전면에 PNG를 깐 형태라 텍스트 편집은 불가하지만
        발표 도구(슬라이드쇼·노트·인쇄)를 그대로 쓸 수 있다."""
        self._require_approved(approval_digest, approval_store)
        self._verify_output_receipt(approval_digest)
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            print("[PPTX] python-pptx 미설치 — pip install python-pptx")
            return None
        out = os.path.join(self.dir, "out")
        fs = sorted(glob.glob(os.path.join(out, "slide_*.png")))
        if not fs:
            print("[PPTX] 슬라이드 없음 — build() 먼저")
            return None
        prs = Presentation()
        prs.slide_width = Inches(13.333)      # 16:9
        prs.slide_height = Inches(7.5)
        prs.core_properties.comments = (
            "Dayoun approval envelope " + self._approval_digest
        )
        blank = prs.slide_layouts[6]
        for f in fs:
            sl = prs.slides.add_slide(blank)
            sl.shapes.add_picture(f, 0, 0, width=prs.slide_width,
                                  height=prs.slide_height)
        p = os.path.join(self.dir, (name or self.title or "deck") + ".pptx")
        prs.save(p)
        print("[PPTX] %s (%.1f MB)" % (p, os.path.getsize(p) / 1024 / 1024))
        return p

    def pdf(self, approval_digest, approval_store, name=None):
        import fitz
        self._require_approved(approval_digest, approval_store)
        self._verify_output_receipt(approval_digest)
        out = os.path.join(self.dir, "out")
        fs = sorted(glob.glob(os.path.join(out, "slide_*.png")))
        doc = fitz.open()
        doc.set_metadata(
            {
                "producer": "dayoun scene-deck",
                "subject": "Approval envelope " + self._approval_digest,
            }
        )
        for f in fs:
            pg = doc.new_page(width=1280, height=720)
            pg.insert_image(pg.rect, filename=f)
        p = os.path.join(self.dir, (name or self.title or "deck") + ".pdf")
        doc.save(p, deflate=True)
        print("[PDF] %s (%.1f MB)" % (p, os.path.getsize(p) / 1024 / 1024))
        return p

    # ══════ 저장/로드 (revise.py 연동) ══════
    def save(self):
        spec = {"domain": self.domain, "foot": self.foot, "title": self.title,
                "family": self.preset["fonts"]["head"], "slides": self.slides}
        if self.brief is not None:
            spec["brief"] = self.brief
        p = os.path.join(self.dir, "spec.json")
        with open(p, "w", encoding="utf-8") as handle:
            json.dump(spec, handle, ensure_ascii=False, indent=1)
        return p

    @classmethod
    def from_brief(cls, brief, domain="it", foot="", out_dir=None):
        """확인된 요구사항 브리프로 덱을 시작한다.

        dict, JSON 경로, 또는 JSON 펜스를 가진 deck_brief.md를 받는다.
        누락 정보가 있거나 사용자 확인 전이면 ``IntakeBlocked``가 다음 질문과
        함께 발생해 조기 생성을 막는다.
        """
        from intake import load_brief, require_confirmed
        if isinstance(brief, (str, os.PathLike)):
            brief = load_brief(brief)
        normalized = require_confirmed(brief)
        title = normalized.get("title") or normalized.get("deck_id") or "deck"
        return cls(domain=domain, foot=foot, title=title, out_dir=out_dir,
                   brief=normalized)

    @classmethod
    def load(cls, spec_path):
        with open(spec_path, encoding="utf-8") as handle:
            spec = json.load(handle)
        d = cls(domain=spec.get("domain", "it"), foot=spec.get("foot", ""),
                title=spec.get("title", ""), out_dir=os.path.dirname(spec_path),
                brief=spec.get("brief"))
        d.slides = spec["slides"]
        return d

    def revise(self):
        """revise.Deck 핸들 반환 — 한 줄 명령으로 수정"""
        from revise import Deck as RDeck
        return RDeck({"domain": self.domain, "foot": self.foot,
                      "family": self.preset["fonts"]["head"], "slides": self.slides},
                     os.path.join(self.dir, "spec.json"))

    def summary(self):
        print("=== %s (%s) ===" % (self.title or "덱", self.preset["key"]))
        for i, s in enumerate(self.slides, 1):
            ex = [k for k in ("num", "chips", "items") if s.get(k)]
            if s.get("rows"):
                ex.append("%d열×%d행" % (len(s["columns"]), len(s["rows"])))
            for k, tag in (("blocks", "대비"), ("bars", "막대")):
                if s.get(k):
                    ex.append("%s%d" % (tag, len(s[k])))
            if s.get("quadrants"):
                ex.append("2×2")
            # 씬이 필요한 구도인데 씬이 없을 때만 경고 — 정보 구도는 원래 씬이 없다
            if (not s.get("scene_body")
                    and s["lay"] not in ("COVER", "CLOSING", "AGENDA",
                                         "TABLE", "EXAMPLE", "MATRIX", "BAR",
                                         "FLOW", "GENEALOGY", "PROMPT")):
                ex.append("씬없음")
            print("  %02d [%s] %-30s %s" % (i, s["lay"], s["head"][0][:30],
                                            "+".join(ex) or ""))


if __name__ == "__main__":
    print(__doc__.split("## 사용법")[1].strip()[:600])
