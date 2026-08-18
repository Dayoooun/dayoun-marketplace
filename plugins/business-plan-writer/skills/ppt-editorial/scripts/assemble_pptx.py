"""Assemble approved rendered PNGs into delivered PPTX and optional PDF.

The approval envelope is resolved by digest from a content-addressed store.
Loose payload, deck-brief, manifest, or HWPX paths are not accepted.
"""
import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from approved_inputs import (
    ApprovalError,
    IMAGE_RENDERER_VERSION,
    resolve_approved_bundle,
    resolve_value,
    verify_image_render_receipt,
)


def main():
    for stream in (sys.stdout, sys.stderr):
        reconfigure = getattr(stream, "reconfigure", None)
        if callable(reconfigure):
            reconfigure(encoding="utf-8", errors="backslashreplace")
    ap = argparse.ArgumentParser()
    ap.add_argument("png_dir")
    ap.add_argument("out")
    ap.add_argument("--approval-digest", required=True)
    ap.add_argument("--approval-store", required=True)
    ap.add_argument("--render-receipt-digest", required=True)
    ap.add_argument("--pdf", action="store_true")
    args = ap.parse_args()

    try:
        store = Path(args.approval_store).resolve()
        envelope = resolve_value(store, args.approval_digest)
        _payload, briefs, _manifest = resolve_approved_bundle(
            store,
            envelope,
            expected_mode="image-first",
            expected_renderer_version=IMAGE_RENDERER_VERSION,
        )
    except ApprovalError as exc:
        raise SystemExit(f"BLOCK: {exc}") from exc

    png_dir = Path(args.png_dir).resolve()
    pngs = sorted(png_dir.glob("*.png"))
    if not pngs:
        raise SystemExit(f"no png in {png_dir}")
    try:
        receipt = resolve_value(store, args.render_receipt_digest)
        verify_image_render_receipt(
            receipt,
            approval_digest=args.approval_digest,
            png_dir=png_dir,
            expected_slide_ids=[
                f"{brief_index}:{slide['slide']}"
                for brief_index, brief in enumerate(briefs, 1)
                for slide in brief.get("slides", [])
            ],
        )
    except ApprovalError as exc:
        raise SystemExit(f"BLOCK: {exc}") from exc
    expected_slides = sum(len(brief.get("slides", [])) for brief in briefs)
    if expected_slides != len(pngs):
        raise SystemExit(
            f"BLOCK: approved briefs require {expected_slides} slides, rendered folder has {len(pngs)}"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.comments = f"Dayoun approval envelope {args.approval_digest}"
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, prs.slide_width, prs.slide_height)
    prs.save(args.out)
    print(f"{len(pngs)} slide(s) -> {args.out}")

    if args.pdf:
        import fitz
        doc = fitz.open()
        doc.set_metadata(
            {
                "producer": "dayoun ppt-editorial",
                "subject": f"Approval envelope {args.approval_digest}",
            }
        )
        for png in pngs:
            img = fitz.open(str(png))
            rect = img[0].rect
            pdf_bytes = img.convert_to_pdf()
            img.close()
            src = fitz.open("pdf", pdf_bytes)
            page = doc.new_page(width=rect.width, height=rect.height)
            page.show_pdf_page(page.rect, src, 0)
        pdf_out = str(Path(args.out).with_suffix(".pdf"))
        # [2026-07-29] 뷰어가 PDF를 열어 두면 저장이 Permission denied로 죽는다.
        # 스택 트레이스 대신 무엇을 해야 하는지 한 줄로 말한다.
        try:
            doc.save(pdf_out)
        except Exception as e:
            if "Permission denied" in str(e) or "denied" in str(e).lower():
                sys.exit(f"PDF를 저장할 수 없다 — {pdf_out} 가 뷰어에 열려 있다. 닫고 다시 실행해라")
            raise
        print(f"pdf -> {pdf_out}")


if __name__ == "__main__":
    main()
