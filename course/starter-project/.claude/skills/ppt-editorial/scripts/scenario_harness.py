#!/usr/bin/env python3
"""Render and validate real-world PPT Editorial scenarios from structured data."""
from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
import shutil
import sys
from pathlib import Path
from typing import Any


from html_slide_renderer import HEIGHT, WIDTH, render_job, validate_graph

def configure_utf8_output(stream) -> None:
    reconfigure = getattr(stream, "reconfigure", None)
    if callable(reconfigure):
        reconfigure(encoding="utf-8")


HERE = Path(__file__).resolve().parent
SKILL_ROOT = HERE.parent
DEFAULT_CATALOG = SKILL_ROOT / "references" / "scenario_catalog.json"
DEFAULT_CONTRACT = SKILL_ROOT / "references" / "render_contract.json"
DEFAULT_FEWSHOTS = SKILL_ROOT / "references" / "scenario_fewshots.json"


class ScenarioHarnessError(RuntimeError):
    pass


def _reject_json_constant(value: str):
    raise ValueError(f"non-finite JSON constant {value!r} is not allowed")


def _non_finite_paths(value: Any, path: str = "$") -> list[str]:
    if isinstance(value, float) and not math.isfinite(value):
        return [path]
    if isinstance(value, dict):
        paths = []
        for key, child in value.items():
            paths.extend(_non_finite_paths(child, f"{path}.{key}"))
        return paths
    if isinstance(value, list):
        paths = []
        for index, child in enumerate(value):
            paths.extend(_non_finite_paths(child, f"{path}[{index}]"))
        return paths
    return []


def _read_json(path: Path) -> dict:
    try:
        data = json.loads(
            path.read_text(encoding="utf-8"),
            parse_constant=_reject_json_constant,
        )
    except (OSError, ValueError) as error:
        raise ScenarioHarnessError(f"cannot read strict JSON {path}: {error}") from error
    if not isinstance(data, dict):
        raise ScenarioHarnessError(f"JSON root must be an object: {path}")
    return data


def _digest(path: Path) -> str:
    return "sha256:" + hashlib.sha256(path.read_bytes()).hexdigest()


def _source_manifest(
    catalog_path: Path,
    contract_path: Path,
    fewshots_path: Path,
    catalog: dict,
) -> dict:
    paths = {
        Path(__file__).resolve(),
        (HERE / "html_slide_renderer.py").resolve(),
        (HERE / "harness_smoke.py").resolve(),
        (HERE / "codex_parallel_gen.py").resolve(),
        (SKILL_ROOT / "SKILL.md").resolve(),
        (SKILL_ROOT / "style_profiles.json").resolve(),
        catalog_path.resolve(),
        contract_path.resolve(),
        fewshots_path.resolve(),
        (SKILL_ROOT / "references" / "image-first-workflow.md").resolve(),
    }
    for scenario in catalog.get("scenarios", []):
        spec = scenario.get("spec", {})
        if spec.get("layout") != "image" or not spec.get("imagePath"):
            continue
        asset = Path(spec["imagePath"])
        if not asset.is_absolute():
            asset = (catalog_path.parent / asset).resolve()
        paths.add(asset)
    rows = [
        {"path": str(path), "digest": _digest(path), "bytes": path.stat().st_size}
        for path in sorted(paths, key=lambda item: str(item).lower())
        if path.is_file()
    ]
    aggregate = hashlib.sha256(
        "".join(row["path"] + row["digest"] for row in rows).encode("utf-8")
    ).hexdigest()
    return {"sourceHash": "sha256:" + aggregate, "files": rows}


def _validate_context(scenario: dict, contract: dict) -> list[str]:
    context = scenario.get("context", {})
    return [
        f"missing context.{key}"
        for key in contract["context"]["required"]
        if not str(context.get(key, "")).strip()
    ]


def _resolve_context_decision(
    scenario: dict,
    contract: dict,
    fewshots: dict,
) -> tuple[dict, list[str]]:
    context = scenario.get("context", {})
    visual_role = str(context.get("visualRole", "")).strip()
    routing = contract.get("contextRouting", {})
    layout = routing.get("visualRoleToLayout", {}).get(visual_role)
    errors: list[str] = []
    if not layout:
        errors.append(f"no context route for visualRole {visual_role!r}")
        return {}, errors
    renderer = routing.get("rendererByLayout", {}).get(layout)
    if not renderer:
        errors.append(f"no renderer route for layout {layout!r}")
    matched = [
        shot["id"]
        for shot in fewshots.get("fewShots", [])
        if shot.get("decision", {}).get("layout") == layout
    ]
    if not matched:
        errors.append(f"no few-shot covers derived layout {layout!r}")
    return {
        "renderer": renderer,
        "layout": layout,
        "routeSource": f"context.visualRole:{visual_role}",
        "fewShotIds": matched,
    }, errors


def _validate_semantics(scenario: dict, contract: dict, decision: dict) -> list[str]:
    spec = scenario.get("spec", {})
    errors: list[str] = []
    layout = spec.get("layout")
    allowed = set(contract["rendererDecision"]["htmlLayouts"])
    if layout not in allowed:
        errors.append(f"layout {layout!r} is not admitted by render contract")
    if decision.get("layout") and layout != decision["layout"]:
        errors.append(
            f"context derives layout {decision['layout']!r}, "
            f"but spec.layout is {layout!r}"
        )
    expected_renderer = contract.get("contextRouting", {}).get(
        "rendererByLayout", {}
    ).get(layout)
    if decision.get("renderer") and expected_renderer != decision["renderer"]:
        errors.append(
            f"context derives renderer {decision['renderer']!r}, "
            f"but contract maps layout to {expected_renderer!r}"
        )
    if layout == "network":
        errors.extend(validate_graph(spec.get("graph", {})))
        for node in spec.get("graph", {}).get("nodes", []):
            if "x" in node or "y" in node:
                errors.append(f"node {node.get('id')!r} hardcodes x/y")
    if layout == "image":
        role = spec.get("assetRole")
        if role not in contract["rendererDecision"]["codexAssetRoles"]:
            errors.append(f"image assetRole {role!r} is not admitted")
        if role == "photo":
            focal = contract.get("photoFocal", {})
            if focal.get("requireImagePosition") and not spec.get("imagePosition"):
                errors.append("photo scenario requires imagePosition")
            minimum_scale = float(focal.get("minScaleWhenReframing", 1.0))
            if float(spec.get("imageScale", 1.0)) < minimum_scale:
                errors.append(
                    f"photo imageScale {spec.get('imageScale', 1.0)} "
                    f"is below {minimum_scale}"
                )
    if layout == "table":
        expected = contract.get("horizontalBalance", {}).get("tableColumns", [])
        actual = []
        for column in spec.get("columns", []):
            raw = str(column.get("width", "")).strip()
            actual.append(round(float(raw.rstrip("%")) / 100, 2) if raw.endswith("%") else None)
        if expected and actual != expected:
            errors.append(f"table column ratios {actual} != {expected}")
    copy = " ".join(
        str(value)
        for key, value in spec.items()
        if key in {"title", "subtitle", "caption", "callout", "note"}
    )
    if "—" in copy:
        errors.append("new slide copy contains forbidden em dash")
    serialized_spec = json.dumps(spec, ensure_ascii=False)
    audience = str(scenario.get("context", {}).get("audience", "")).lower()
    if (
        "public" in audience
        and re.search(r"010-\d{4}-\d{4}", serialized_spec)
        and not any(
            marker in serialized_spec
            for marker in ("테스트 전용", "가상사용자")
        )
    ):
        errors.append("public scenario contains realistic unmarked phone data")
    return errors


def _resolve_assets(spec: dict, catalog_dir: Path) -> dict:
    resolved = json.loads(json.dumps(spec, ensure_ascii=False, allow_nan=False))
    if resolved.get("layout") == "image" and resolved.get("imagePath"):
        path = Path(resolved["imagePath"])
        if not path.is_absolute():
            path = (catalog_dir / path).resolve()
        resolved["imagePath"] = str(path)
    return resolved


def _check_non_uniform(path: Path) -> bool:
    import numpy as np
    from PIL import Image
    with Image.open(path) as image:
        array = np.asarray(image.convert("RGB"), dtype=np.float32)
    return float(array.std()) >= 4.0


def _validate_receipt(scenario: dict, receipt: dict, contract: dict) -> list[str]:
    errors: list[str] = []
    non_finite = _non_finite_paths(receipt)
    if non_finite:
        errors.append(
            "receipt contains non-finite values at " + ", ".join(non_finite)
        )
    gate = contract["qualityGate"]
    if receipt.get("size") != gate["size"]:
        errors.append(f"size mismatch: {receipt.get('size')} != {gate['size']}")
    if len(receipt.get("overflow", [])) > gate["maxOverflowElements"]:
        errors.append(f"overflow elements: {receipt['overflow']}")
    renderer_digest = _digest(HERE / "html_slide_renderer.py")
    if receipt.get("rendererDigest") != renderer_digest:
        errors.append("renderer digest mismatch")
    if not str(receipt.get("specDigest", "")).startswith("sha256:"):
        errors.append("spec digest missing")

    layout = scenario["spec"]["layout"]
    if layout not in contract["verticalBalance"]["exceptions"]:
        content = receipt.get("content") or {}
        body = receipt.get("meaningfulBody") or {}
        top_min, top_max = contract["verticalBalance"]["structuralBodyTop"]
        bottom_min, bottom_max = contract["verticalBalance"]["meaningfulBodyBottom"]
        top = content.get("top")
        bottom = body.get("bottom")
        if top is None or not top_min <= top <= top_max:
            errors.append(f"content top {top} outside [{top_min}, {top_max}]")
        if bottom is None or not bottom_min <= bottom <= bottom_max:
            errors.append(f"meaningful body bottom {bottom} outside [{bottom_min}, {bottom_max}]")

    if layout == "modules":
        balance = receipt.get("moduleBalance") or {}
        limits = contract["moduleFeatureBalance"]["maxChildMargin"]
        for side, limit in limits.items():
            value = balance.get(side)
            if value is None or value > limit:
                errors.append(f"module feature {side} margin {value} exceeds {limit}")

    if layout == "network":
        expected_nodes = len(scenario["spec"]["graph"]["nodes"])
        expected_edges = len(scenario["spec"]["graph"]["edges"])
        graph = receipt.get("graph", {})
        if graph.get("nodes") != expected_nodes:
            errors.append(f"graph node receipt {graph.get('nodes')} != {expected_nodes}")
        if graph.get("edges") != expected_edges:
            errors.append(f"graph edge receipt {graph.get('edges')} != {expected_edges}")
        if graph.get("edgeLabels") != expected_edges:
            errors.append(f"graph edge labels {graph.get('edgeLabels')} != {expected_edges}")
        visibility = graph.get("visibility", [])
        if len(visibility) != expected_edges:
            errors.append(f"graph visibility receipts {len(visibility)} != {expected_edges}")
        else:
            for expected, visible in zip(
                scenario["spec"]["graph"]["edges"],
                visibility,
                strict=True,
            ):
                for field in ("source", "target", "direction", "label"):
                    if str(visible.get(field)) != str(expected.get(field)):
                        errors.append(
                            f"graph visibility {field} {visible.get(field)!r} "
                            f"!= {expected.get(field)!r}"
                        )
                if not visible.get("hasMarkerEnd"):
                    errors.append(
                        f"edge {expected['source']}->{expected['target']} lacks end marker"
                    )
                if expected.get("direction") == "bidirectional" and not visible.get("hasMarkerStart"):
                    errors.append(
                        f"edge {expected['source']}->{expected['target']} lacks start marker"
                    )
                for field in (
                    "startOccluded",
                    "endOccluded",
                    "pathOccluded",
                    "labelOccluded",
                ):
                    if visible.get(field):
                        errors.append(
                            f"edge {expected['source']}->{expected['target']} "
                            f"{field}"
                        )
                if float(visible.get("length", 0)) <= 20:
                    errors.append(
                        f"edge {expected['source']}->{expected['target']} is too short"
                    )

    if layout == "process":
        gap = receipt.get("internalGaps", {}).get("processToNote")
        maximum = contract.get("internalGap", {}).get("processToNoteMax")
        if maximum is not None:
            finite = isinstance(gap, (int, float)) and math.isfinite(float(gap))
            if not finite or gap > maximum:
                errors.append(f"process-to-note gap {gap} exceeds {maximum}")

    if layout == "image":
        source = receipt.get("sourceAsset")
        digest = receipt.get("sourceAssetDigest")
        if not source:
            errors.append("image source asset missing from receipt")
        elif digest != _digest(Path(source)):
            errors.append("image source asset digest mismatch")
    return errors


def _contact_sheet(images: list[tuple[str, Path]], out: Path) -> None:
    from PIL import Image, ImageDraw
    thumb_w, thumb_h = 668, 376
    gap, label_h = 24, 30
    columns = 2
    rows = (len(images) + columns - 1) // columns
    sheet = Image.new("RGB", (columns * thumb_w + (columns + 1) * gap, rows * (thumb_h + label_h) + (rows + 1) * gap), "#151515")
    draw = ImageDraw.Draw(sheet)
    for index, (label, path) in enumerate(images):
        row, column = divmod(index, columns)
        x = gap + column * (thumb_w + gap)
        y = gap + row * (thumb_h + label_h + gap)
        with Image.open(path) as image:
            thumb = image.convert("RGB").resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        sheet.paste(thumb, (x, y))
        draw.text((x, y + thumb_h + 5), label, fill="white")
    out.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out)


def _assemble(images: list[tuple[str, Path]], pptx_path: Path, pdf_path: Path) -> None:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for _, path in images:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, presentation.slide_width, presentation.slide_height)
    presentation.save(pptx_path)

    pdf = fitz.open()
    for _, path in images:
        page = pdf.new_page(width=960, height=540)
        page.insert_image(page.rect, filename=str(path))
    pdf.save(pdf_path)
    pdf.close()


def run(catalog_path: Path, contract_path: Path, fewshots_path: Path, out_dir: Path) -> dict:
    catalog = _read_json(catalog_path)
    contract = _read_json(contract_path)
    fewshots = _read_json(fewshots_path)
    scenarios = catalog.get("scenarios", [])
    if not scenarios:
        raise ScenarioHarnessError("scenario catalog is empty")
    if not fewshots.get("fewShots"):
        raise ScenarioHarnessError("few-shot catalog is empty")
    source_manifest = _source_manifest(
        catalog_path,
        contract_path,
        fewshots_path,
        catalog,
    )

    out_dir.mkdir(parents=True, exist_ok=True)
    slides_dir = out_dir / "slides"
    if slides_dir.exists():
        shutil.rmtree(slides_dir)
    slides_dir.mkdir(parents=True)

    results = []
    images: list[tuple[str, Path]] = []
    for index, scenario in enumerate(scenarios, 1):
        scenario_id = str(scenario.get("id", f"S{index:03d}"))
        errors = _validate_context(scenario, contract)
        decision, decision_errors = _resolve_context_decision(
            scenario,
            contract,
            fewshots,
        )
        errors.extend(decision_errors)
        errors.extend(_validate_semantics(scenario, contract, decision))
        spec = _resolve_assets(scenario.get("spec", {}), catalog_path.parent)
        out = slides_dir / f"slide_{index:02d}_{scenario_id}.png"
        if not errors:
            try:
                render_job({"renderer": "html", "htmlSpec": spec, "out": str(out)}, out_dir)
            except Exception as error:
                errors.append(f"render failed: {type(error).__name__}: {error}")
        receipt_path = out.with_suffix(".layout.json")
        receipt = {}
        if out.is_file() and receipt_path.is_file():
            receipt = _read_json(receipt_path)
            errors.extend(_validate_receipt(scenario, receipt, contract))
            if gate := contract.get("qualityGate"):
                if gate.get("requireNonUniformImage") and not _check_non_uniform(out):
                    errors.append("rendered image is visually uniform")
            images.append((scenario_id, out))
        else:
            errors.append("render or layout receipt missing")
        results.append({
            "id": scenario_id,
            "layout": spec.get("layout"),
            "decision": decision,
            "status": "PASS" if not errors else "FAIL",
            "errors": errors,
            "image": str(out) if out.is_file() else None,
            "receipt": str(receipt_path) if receipt_path.is_file() else None,
            "imageDigest": _digest(out) if out.is_file() else None,
            "receiptDigest": _digest(receipt_path) if receipt_path.is_file() else None,
        })

    contact = out_dir / "contact-sheet.png"
    pptx = out_dir / "scenario-harness.pptx"
    pdf = out_dir / "scenario-harness.pdf"
    artifact_errors: list[str] = []
    if images:
        _contact_sheet(images, contact)
        _assemble(images, pptx, pdf)
    quality_gate = contract.get("qualityGate", {})
    if quality_gate.get("requireSourceHash") and not source_manifest["sourceHash"].startswith("sha256:"):
        artifact_errors.append("source hash missing")
    if quality_gate.get("requireInputOutputDigests"):
        for result in results:
            if not result.get("imageDigest") or not result.get("receiptDigest"):
                artifact_errors.append(
                    f"output digest missing for {result['id']}"
                )
    if quality_gate.get("requireContactSheet"):
        if not contact.is_file() or not _check_non_uniform(contact):
            artifact_errors.append("contact sheet missing or visually uniform")
    import fitz
    from pptx import Presentation
    pptx_pages = len(Presentation(pptx).slides) if pptx.is_file() else 0
    if quality_gate.get("requirePptxPagesMatch") and pptx_pages != len(scenarios):
        artifact_errors.append(
            f"PPTX pages {pptx_pages} != scenario count {len(scenarios)}"
        )
    pdf_document = fitz.open(pdf) if pdf.is_file() else None
    pdf_pages = len(pdf_document) if pdf_document is not None else 0
    if pdf_document is not None:
        pdf_document.close()
    if quality_gate.get("requirePdfPagesMatch") and pdf_pages != len(scenarios):
        artifact_errors.append(
            f"PDF pages {pdf_pages} != scenario count {len(scenarios)}"
        )
    if len(images) != len(scenarios):
        artifact_errors.append(
            f"rendered images {len(images)} != scenario count {len(scenarios)}"
        )
    scenario_failures = sum(result["status"] == "FAIL" for result in results)
    overall_status = (
        "PASS"
        if scenario_failures == 0 and not artifact_errors
        else "FAIL"
    )
    report = {
        "schemaVersion": 1,
        "status": overall_status,
        "sourceHash": source_manifest["sourceHash"],
        "sourceManifest": source_manifest,
        "catalog": {"path": str(catalog_path), "digest": _digest(catalog_path)},
        "contract": {"path": str(contract_path), "digest": _digest(contract_path)},
        "fewShots": {"path": str(fewshots_path), "digest": _digest(fewshots_path), "count": len(fewshots["fewShots"])},
        "scenarioCount": len(scenarios),
        "renderedCount": len(images),
        "passed": sum(result["status"] == "PASS" for result in results),
        "failed": scenario_failures,
        "artifactFailed": len(artifact_errors),
        "artifactErrors": artifact_errors,
        "results": results,
        "artifacts": {
            "contactSheet": {
                "path": str(contact),
                "digest": _digest(contact) if contact.is_file() else None,
            },
            "pptx": {
                "path": str(pptx),
                "digest": _digest(pptx) if pptx.is_file() else None,
                "pages": pptx_pages,
            },
            "pdf": {
                "path": str(pdf),
                "digest": _digest(pdf) if pdf.is_file() else None,
                "pages": pdf_pages,
            },
        },
    }
    report_path = out_dir / "scenario-report.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2, allow_nan=False) + "\n", encoding="utf-8")
    return report


def main() -> int:
    configure_utf8_output(sys.stdout)
    parser = argparse.ArgumentParser(description="Render and validate structured PPT Editorial scenarios")
    parser.add_argument("--catalog", type=Path, default=DEFAULT_CATALOG)
    parser.add_argument("--contract", type=Path, default=DEFAULT_CONTRACT)
    parser.add_argument("--fewshots", type=Path, default=DEFAULT_FEWSHOTS)
    parser.add_argument("--out-dir", type=Path, required=True)
    args = parser.parse_args()
    report = run(args.catalog.resolve(), args.contract.resolve(), args.fewshots.resolve(), args.out_dir.resolve())
    print(json.dumps({key: report[key] for key in ("status", "scenarioCount", "renderedCount", "passed", "failed", "artifactFailed")}, ensure_ascii=False, allow_nan=False))
    return 0 if report["status"] == "PASS" else 1


if __name__ == "__main__":
    raise SystemExit(main())
